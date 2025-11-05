import os
import io
from fastapi import APIRouter,UploadFile,File,Form
import tempfile
from memory.memory_client import *
# --- Core Libraries ---
import fitz  # PyMuPDF for PDFs
import docx  # python-docx for .docx
from pptx import Presentation  # python-pptx for .pptx
from PIL import Image  # Pillow for handling images

file_router = APIRouter()
# --- OCR Library ---
try:
    import pytesseract
    import shutil
    tesseract_cmd = shutil.which('tesseract')
    if not tesseract_cmd:
        possible_paths = [
            r'C:\Program Files\Tesseract-OCR\tesseract.exe',
            r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
            r'C:\Users\AppData\Local\Programs\Tesseract-OCR\tesseract.exe'
        ]
        
        for path in possible_paths:
            if os.path.exists(path):
                tesseract_cmd = path
                break
    
    if tesseract_cmd:
        pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
        print(f"Found Tesseract at: {tesseract_cmd}")
        try:
            version = pytesseract.get_tesseract_version()
            print(f"Tesseract version: {version}")
        except Exception as e:
            print(f"Error checking Tesseract version: {e}")
            pytesseract = None
    else:
        print("Error: Tesseract executable not found!")
        print("Please ensure Tesseract is installed and in one of these locations:")
        for path in possible_paths:
            print(f"- {path}")
        pytesseract = None

except ImportError:
    print("Warning: pytesseract not installed. OCR functionality will not be available.")
    print("Install it with: pip install pytesseract")
    pytesseract = None

def _ocr_image(image: Image.Image) -> str:
    """Helper function to perform OCR on a PIL Image object."""
    if pytesseract is None:
        return "[OCR Skipped: pytesseract not installed]"
    try:
        # Perform OCR
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        return f"[OCR Error: {e}]"

def _extract_text_from_pdf(file_path: str) -> str:
    """Extracts text from a PDF, attempting direct text extraction first, then OCR as a fallback."""
    all_text = []
    try:
        doc = fitz.open(file_path)
        
        # 1. Attempt direct text extraction
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            all_text.append(page.get_text())
            
        full_text = "\n".join(all_text).strip()
        
        # 2. Check if text is sparse (likely a scanned PDF)
        # We set a threshold (e.g., less than 100 non-whitespace chars) to trigger OCR.
        if len(full_text) < 100 and pytesseract is not None:
            all_text = [] # Reset to store OCR results
            print(f"Info: '{os.path.basename(file_path)}' appears to be scanned. Falling back to OCR.")
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Render page to a high-DPI image (pixmap)
                # 300 DPI is good for OCR
                pix = page.get_pixmap(dpi=300)
                
                # Convert pixmap to a PIL Image
                img_bytes = pix.tobytes("png")
                img = Image.open(io.BytesIO(img_bytes))
                
                # Perform OCR on the page image
                page_text = _ocr_image(img)
                all_text.append(page_text)
        
        doc.close()
        
    except Exception as e:
        return f"[PDF Error: {e}]"
        
    return "\n".join(all_text)

def _extract_text_from_docx(file_path: str) -> str:
    """Extracts text from a .docx file."""
    try:
        doc = docx.Document(file_path)
        all_text = [para.text for para in doc.paragraphs]
        return "\n".join(all_text)
    except Exception as e:
        return f"[DOCX Error: {e}]"

def _extract_text_from_pptx(file_path: str) -> str:
    """Extracts text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        all_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    all_text.append(shape.text)
        return "\n".join(all_text)
    except Exception as e:
        return f"[PPTX Error: {e}]"

def _extract_text_from_txt(file_path: str) -> str:
    """Extracts text from a .txt file."""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"[TXT Error: {e}]"

def _extract_text_from_image(file_path: str) -> str:
    """Extracts text from an image file using OCR."""
    if pytesseract is None:
        return "[OCR Error: pytesseract not installed or Tesseract engine not found]"
    try:
        img = Image.open(file_path)
        return _ocr_image(img)
    except Exception as e:
        return f"[Image Error: {e}]"



def extract_text(file_path: str,ext:str) -> str:
    if not os.path.exists(file_path):
        return f"Error: File not found at {file_path}"

    try:
        if ext=="pdf":
            raw_text = _extract_text_from_pdf(file_path)
            
        elif ext=="ocx":
            raw_text = _extract_text_from_docx(file_path)
            
        elif ext=="ptx":
            raw_text = _extract_text_from_pptx(file_path)
            
        elif ext=="txt":
            raw_text = _extract_text_from_txt(file_path)
            
        elif ext in ["png", "jpg", "peg", "bmp", "iff"]:
            raw_text = _extract_text_from_image(file_path)
            
        else:
            return f"Error: Unsupported file type"

    except Exception as e:
        return f"Error: An unexpected error occurred while processing {file_path}: {e}"

    cleaned_text = " ".join(raw_text.split())
    
    return cleaned_text


@file_router.post('/api/memory/uploadfile')
async def query_file(uploaded_file:UploadFile=File(...)):
    with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
        tmp_file.write(uploaded_file.file.read())
        file_path = tmp_file.name
    print("FILE_PATH: ",file_path)
    file_ext=uploaded_file.filename[-3:]
    content=extract_text(file_path=file_path,ext=file_ext)
    print("FILE EXT=",file_ext)
    result= await add_single_memory(context=content,user_id="central-memories")
    return {"status":"success"}